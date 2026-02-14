from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List


Slide = Dict[str, Any]


class SlideGenerator:
    """
    Generate markdown slides (and optionally pptx) from a slide plan.
    """

    def to_markdown(self, slides: List[Slide]) -> str:
        """
        Convert a list of slide dicts into markdown slides separated by `---`.
        """
        parts: list[str] = []
        for idx, slide in enumerate(slides):
            title = str(slide.get("title", "")).strip() or f"Slide {idx + 1}"
            bullets = slide.get("bullets", [])
            if not isinstance(bullets, list):
                bullets = [str(bullets)]

            lines: list[str] = []
            # Use level-1 heading per slide for good compatibility with Marp / reveal.js.
            lines.append(f"# {title}")
            lines.append("")
            for b in bullets:
                text = str(b).strip()
                if not text:
                    continue
                lines.append(f"- {text}")
            parts.append("\n".join(lines).strip())

        return "\n\n---\n\n".join(parts) + "\n"

    def to_markdown_with_front_matter(
        self,
        slides: List[Slide],
        *,
        title: str = "Repo2Slides",
        author: str = "",
        theme: str = "default",
        engine: str = "marp",
    ) -> str:
        """
        Markdown slides with YAML front matter.

        - engine="marp": add Marp-compatible fields (marp:true, paginate:true)
        - engine="plain": generic YAML only
        """
        meta_lines = ["---"]
        if engine == "marp":
            meta_lines.append("marp: true")
            meta_lines.append("paginate: true")
        meta_lines.append(f'title: "{title}"')
        if author:
            meta_lines.append(f'author: "{author}"')
        if theme:
            meta_lines.append(f'theme: "{theme}"')
        meta_lines.append(f'generated_at: "{datetime.utcnow().isoformat(timespec="seconds")}Z"')
        meta_lines.append("---")
        meta = "\n".join(meta_lines) + "\n\n"
        return meta + self.to_markdown(slides)

    def to_pptx(self, slides: List[Slide], output_path: Path) -> None:
        """
        Generate a basic PPTX file using python-pptx.
        """
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as e:  # pragma: no cover - import guard
            raise ImportError(
                "python-pptx is required for PPTX export. Install via 'pip install python-pptx'."
            ) from e

        prs = Presentation()

        # Remove default slides
        while len(prs.slides) > 0:
            r = prs.slides._sldIdLst[0]
            prs.slides._sldIdLst.remove(r)

        for idx, slide_data in enumerate(slides):
            layout = prs.slide_layouts[1] if idx > 0 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)

            title = str(slide_data.get("title", "")).strip() or f"Slide {idx + 1}"
            bullets = slide_data.get("bullets", [])
            if not isinstance(bullets, list):
                bullets = [str(bullets)]

            # Title
            if slide.shapes.title is not None:
                slide.shapes.title.text = title

            # Content placeholder
            body = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # BODY
                    body = shape
                    break
            if body is not None:
                tf = body.text_frame
                tf.clear()
                first = True
                for b in bullets:
                    text = str(b).strip()
                    if not text:
                        continue
                    if first:
                        tf.text = text
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = text
                        p.level = 0

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

